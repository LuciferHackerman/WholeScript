import xlsxwriter
from docx import Document
from docx.shared import Pt, RGBColor
from pptx import Presentation

# !!!EXCEL!!!
# пример данных
chart_data = [
    {'name': 'Lorem', 'value': 23},
    {'name': 'Ipsum', 'value': 48},
    {'name': 'Dolor', 'value': 15},
    {'name': 'Sit', 'value': 8},
    {'name': 'Amet', 'value': 32},
]
# путь к xlsx файлу
xls_file = 'chart.xlsx'

# рабочаяя книга
workbook = xlsxwriter.Workbook(xls_file)

# добавить новый лист в книгу
worksheet = workbook.add_worksheet()

row_ = 0
col_ = 0

# написать заголовки
worksheet.write(row_, col_, 'NAME')
col_ += 1
worksheet.write(row_, col_, 'VALUE')
row_ += 1

# вписать пример данных
for item in chart_data:
    col_ = 0
    worksheet.write(row_, col_, item['name'])
    col_ += 1
    worksheet.write(row_, col_, item['value'])
    row_ += 1

# создать круговую диаграмму
pie_chart = workbook.add_chart({'type': 'pie'})

# добавить ряды
pie_chart.add_series({
    'name': 'Series Name',
    'categorie': '=Sheet1!$A$3:$A$%s' % row_,
    'values': '=Sheet1!$B$3:$B$%s' % row_,
    'marker': {'type': 'circle'}
})

# вставить круговую диаграмму
worksheet.insert_chart('D2', pie_chart)

# добавить столбцовую диаграмму
column_chart = workbook.add_chart({'type': 'column'})

# добавить ряд к диаграмме
column_chart.add_series({
    'name': 'Series Name',
    'categorie': '=Sheet1!$A$3:$A$%s' % row_,
    'values': '=Sheet1!$B$3:$B$%s' % row_,
    'marker': {'type': 'circle'}
})

# вставить столбцовую диаграмму
worksheet.insert_chart('D2', column_chart)

workbook.close()

# !!!PRESENTATION!!!
prs = Presentation()

# first slide
first_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(first_slide_layout)
title = slide.shapes.title
title.text = "Устинов Егор Максимович"

# second slide
second_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(second_slide_layout)
title = slide.shapes.title
title.text = "Операционные системы"

# third slide
third_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(third_slide_layout)
title = slide.shapes.title
title.text = "Лабораторная работа №2"

prs.save('presentation.pptx')

# !!!DOC FILE!!!
doc = Document()

parOne = doc.add_paragraph()
p = parOne.add_run('Устинов Егор Максимович')
p.font.name = 'Comic Sans MS'
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(0x00, 0xff, 0x00)
doc.add_page_break()

parTwo = doc.add_paragraph()
p = parTwo.add_run('Операционные Системы')
p.font.name = 'Comic Sans MS'
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(0x40, 0xEB, 0x34)
doc.add_page_break()

parOne = doc.add_paragraph()
p = parOne.add_run('Лабораторная работа №2')
p.font.name = 'Comic Sans MS'
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(0x40, 0xEB, 0x34)

doc.save('document.docx')
